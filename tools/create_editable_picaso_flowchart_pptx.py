from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("/Users/xin/Downloads/Daniel_Huang_PICASO_Workflow_Flowchart_EDITABLE.pptx")


def rgb(value: str) -> RGBColor:
    value = value.strip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


INK = "#2b2f33"
DARK = "#62696d"
MID = "#a8afb2"
LINE = "#4b5257"
PAPER = "#f7f8f7"
BLUE = "#c8dce4"
GREEN = "#dbe8df"
ROSE = "#efd4cc"
SAND = "#eee6d6"
GRID_GREEN = "#a2ada7"
GRID_ROSE = "#b5a9a5"


def set_text(shape, text, *, size=10, color=INK, bold=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.07)
    tf.margin_right = Inches(0.07)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    for para in tf.paragraphs:
        para.space_after = Pt(0)
        para.space_before = Pt(0)
        para.line_spacing = 0.92


def add_rect(slide, x, y, w, h, *, fill, line=LINE, radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(1.2)
    return shape


def add_box(slide, x, y, w, h, title, body, *, fill, header=DARK, body_size=8.2):
    box = add_rect(slide, x, y, w, h, fill=fill)
    header_shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x + 0.08),
        Inches(y + 0.08),
        Inches(w - 0.16),
        Inches(0.30),
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = rgb(header)
    header_shape.line.color.rgb = rgb(header)
    set_text(header_shape, title, size=10.5, color="#ffffff", bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    body_shape = slide.shapes.add_textbox(Inches(x + 0.23), Inches(y + 0.52), Inches(w - 0.45), Inches(h - 0.60))
    set_text(body_shape, body, size=body_size, color=INK, valign=MSO_ANCHOR.TOP)
    return box


def add_bar(slide, x, y, w, h, text, *, fill=DARK, size=10):
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(fill)
    bar.line.color.rgb = rgb(LINE)
    bar.line.width = Pt(0.8)
    set_text(bar, text, size=size, color="#ffffff", bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    return bar


def add_grid_item(slide, x, y, w, h, text, *, fill=MID, size=6.9):
    item = add_rect(slide, x, y, w, h, fill=fill)
    set_text(item, text, size=size, color="#ffffff", align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    return item


def add_text(slide, x, y, w, h, text, *, size=12, color=INK, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text(tb, text, size=size, color=color, bold=bold, align=align, valign=MSO_ANCHOR.MIDDLE)
    return tb


def add_arrow(slide, x1, y1, x2, y2, *, width=1.1, dashed=False):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = rgb(LINE)
    line.line.width = Pt(width)
    line.line.end_arrowhead = True
    if dashed:
        line.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return line


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb(PAPER)

    add_text(slide, 1.05, 0.12, 11.25, 0.35, "PICASO Sub-Neptune Planet-Typing Workflow", size=26, bold=True, align=PP_ALIGN.CENTER)
    add_text(
        slide,
        1.55,
        0.48,
        10.25,
        0.22,
        "Aurora schematic based on Daniel Huang's AABC proposal and the timestep/RoadRunner PICASO workflow",
        size=10.5,
        color="#56616a",
        align=PP_ALIGN.CENTER,
    )

    add_text(slide, 0.62, 0.78, 5.4, 0.30, "DIAGRAM OF HOW MODELS ARE COMPUTED", size=14, bold=True)
    add_box(
        slide,
        0.55,
        1.10,
        2.90,
        1.82,
        "INPUTS",
        "Science question: can cloudy, volatile-rich sub-Neptunes near the radius valley look terrestrial in reflected light?\n"
        "Proposal grid: FGK hosts, HZ separations, enrichment up to ~1000x solar, cloud levers, and mission-relevant wavelength windows.\n"
        "Current timestep inputs: PICASO reference data plus SLGRID pressure-temperature and cloud files.",
        fill=BLUE,
        body_size=6.9,
    )
    add_box(
        slide,
        3.82,
        1.25,
        2.25,
        1.06,
        "MODEL SETUP",
        "SystemParams builds each case: Teff, logg, Rp, a_AU, phase_deg, star T/R.\n"
        "PICASO opacity connection sets the wavelength range.",
        fill="#dfe8eb",
        body_size=6.9,
    )
    add_box(
        slide,
        3.82,
        2.48,
        2.25,
        0.96,
        "ATMOSPHERE",
        "resolve_slgrid_files selects paired PT and cloud files.\n"
        "Cloud-column names are normalized before case.clouds(...).",
        fill="#dfe8eb",
        body_size=6.9,
    )
    add_box(
        slide,
        6.48,
        1.25,
        2.62,
        1.06,
        "REFLECTED-LIGHT STEP",
        "case.star + case.gravity + case.atmosphere + case.clouds.\n"
        "case.phase_angle(phase_deg) then spectrum(..., calculation='reflected').",
        fill=GREEN,
        body_size=6.8,
    )
    add_box(
        slide,
        6.48,
        2.40,
        2.62,
        1.08,
        "THERMAL / ANALOG STEP",
        "PICASO thermal uses phase=0, or the hybrid workflow injects EGP g31 IRflux.\n"
        "Earth benchmark spectra and abiotic terrestrial analogs form the comparison set.",
        fill=SAND,
        body_size=6.8,
    )
    add_box(
        slide,
        9.45,
        1.48,
        2.55,
        1.55,
        "POST-PROCESS",
        "Convert wavenumber to wavelength.\n"
        "Use fp/fs or albedo to recover absolute reflected planet flux.\n"
        "Interpolate onto LAM_GRID; integrate Roman CGI bands and proposed HWO windows.\n"
        "Compute f_reflect and spectral similarity metrics.",
        fill=ROSE,
        body_size=6.7,
    )
    add_box(
        slide,
        4.80,
        3.53,
        3.95,
        0.60,
        "FINAL OUTPUT",
        "Sub-Neptune spectral library, confusion-region maps, diagnostic wavelength windows, and planet-typing guidance.",
        fill="#cdd8dc",
        body_size=7.0,
    )
    add_grid_item(slide, 0.70, 2.98, 2.55, 0.40, "Implementation path:\ntimestep/2_9_2026 -> phase-60 hybrid notebook -> migrated picaso4 workflow", fill="#8b9498", size=6.6)

    add_arrow(slide, 3.45, 1.90, 3.82, 1.77)
    add_arrow(slide, 3.45, 2.40, 3.82, 2.95)
    add_arrow(slide, 6.07, 1.78, 6.48, 1.78)
    add_arrow(slide, 6.07, 2.96, 6.48, 2.94)
    add_arrow(slide, 9.10, 1.78, 9.45, 2.05)
    add_arrow(slide, 9.10, 2.92, 9.45, 2.38)
    add_arrow(slide, 10.30, 3.03, 8.35, 3.53)
    add_arrow(slide, 7.75, 3.48, 7.60, 3.53)

    add_text(slide, 0.62, 4.23, 4.6, 0.30, "DIAGRAM OF GRID STRUCTURE", size=14, bold=True)
    # Grid frame.
    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.78), Inches(4.58), Inches(11.78), Inches(2.46))
    frame.fill.background()
    frame.line.color.rgb = rgb(LINE)
    frame.line.width = Pt(1.1)

    add_bar(slide, 0.88, 4.68, 11.58, 0.20, "MODEL", size=10.5)
    rows = [
        ("HOST STAR", 4.96, ["FGK host bins", "stellar spectra", "HZ separation / isolation"], "#9fa7ab"),
        ("PLANET", 5.58, ["radius-valley target", "sub-Neptune envelope", "Teff, logg, Rp, a_AU"], "#a7acae"),
        ("ATMOSPHERE + CLOUDS", 6.12, ["metallicity / MMW", "SLGRID PT profile", "cloud top + opacity", "fsed / cloud fraction"], GRID_GREEN),
        ("OBSERVATION + METRICS", 6.72, ["phase geometry", "0.3-1.0 um current grid; 0.3-2.5 um proposal", "Roman CGI / HWO windows", "band depth + slopes", "CIA / continuum", "noise similarity"], GRID_ROSE),
    ]

    add_bar(slide, 1.10, 4.98, 11.00, 0.18, "HOST STAR", fill="#6f777a", size=9.5)
    x0, gap = 1.35, 0.15
    w3 = (10.25 - 2 * gap) / 3
    for i, text in enumerate(rows[0][2]):
        add_grid_item(slide, x0 + i * (w3 + gap), 5.20, w3, 0.36, text, fill=rows[0][3], size=6.6)
        add_arrow(slide, x0 + i * (w3 + gap) + w3 / 2, 5.56, x0 + i * (w3 + gap) + w3 / 2, 5.66, width=0.7)

    add_bar(slide, 1.10, 5.68, 11.00, 0.18, "PLANET", fill="#6f777a", size=9.5)
    for i, text in enumerate(rows[1][2]):
        add_grid_item(slide, x0 + i * (w3 + gap), 5.90, w3, 0.36, text, fill=rows[1][3], size=6.6)
        add_arrow(slide, x0 + i * (w3 + gap) + w3 / 2, 6.26, x0 + i * (w3 + gap) + w3 / 2, 6.36, width=0.7)

    add_bar(slide, 1.10, 6.38, 11.00, 0.18, "ATMOSPHERE + CLOUDS", fill="#6f777a", size=9.5)
    w4 = (10.25 - 3 * gap) / 4
    for i, text in enumerate(rows[2][2]):
        add_grid_item(slide, x0 + i * (w4 + gap), 6.58, w4, 0.34, text, fill=rows[2][3], size=6.4)
        if i in (1, 2):
            add_arrow(slide, x0 + i * (w4 + gap) + w4 / 2, 6.92, x0 + i * (w4 + gap) + w4 / 2, 7.01, width=0.7)

    add_bar(slide, 1.10, 7.03, 11.00, 0.18, "OBSERVATION + METRICS", fill="#6f777a", size=9.5)
    w6 = (10.25 - 5 * gap) / 6
    for i, text in enumerate(rows[3][2]):
        add_grid_item(slide, x0 + i * (w6 + gap), 7.25, w6, 0.34, text, fill=rows[3][3], size=5.9)

    add_text(
        slide,
        1.75,
        7.25,
        9.8,
        0.18,
        "Fact-check notes: proposal target range is 0.3-2.5 um; current RoadRunner LAM_GRID is 0.3-1.0 um. Current code uses Roman CGI bands; HWO windows are the proposal deliverable.",
        size=5.6,
        color="#5a646b",
        align=PP_ALIGN.CENTER,
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
